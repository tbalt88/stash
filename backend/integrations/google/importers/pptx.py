"""PPTX importer: Drive → fixed-aspect HTML slide page.

Flow:
  1. Drive download (alt=media) → bytes
  2. LibreOffice headless: pptx → pdf
  3. pdftoppm: pdf → one PNG per slide
  4. python-pptx: extract per-slide text (used as a hidden overlay so
     search and AI features can still see the slide contents)
  5. Build HTML — one `<section class="slide">` per slide, each with a
     full-bleed background image (embedded as a base64 data URI to
     avoid auth dance for iframe-loaded images) and an invisible
     text layer for indexing.

Images are inlined to keep the page self-contained. A 20-slide deck
typically lands in the 6–12 MB range; documented trade-off vs. an
S3-hosted asset chain that would need refreshable URLs.
"""

from __future__ import annotations

import asyncio
import base64
import html as html_mod
import logging
import shutil
import tempfile
from pathlib import Path
from uuid import UUID

import asyncpg
import httpx
from pptx import Presentation

from ....database import get_pool

logger = logging.getLogger(__name__)

DRIVE_DOWNLOAD_URL = (
    "https://www.googleapis.com/drive/v3/files/{file_id}?alt=media&supportsAllDrives=true"
)
PDFTOPPM_DPI = 144
LIBREOFFICE_TIMEOUT = 180  # seconds


async def import_pptx_from_drive(
    client: httpx.AsyncClient,
    workspace_id: UUID,
    folder_id: UUID | None,
    user_id: UUID,
    file_id: str,
    name: str,
) -> dict:
    raw = await _drive_download(client, file_id)
    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        pptx_path = td_path / "deck.pptx"
        pptx_path.write_bytes(raw)

        pdf_path = await _convert_to_pdf(pptx_path, td_path)
        png_paths = await _render_pdf_to_pngs(pdf_path, td_path)
        slide_texts = _extract_slide_text(pptx_path)

    html = _build_slide_html(name, png_paths_to_data_uris(png_paths), slide_texts)

    pool = get_pool()
    try:
        row = await pool.fetchrow(
            """
            INSERT INTO pages (
                workspace_id, folder_id, name, content_html,
                content_type, html_layout, created_by
            ) VALUES ($1, $2, $3, $4, 'html', 'fixed-aspect', $5)
            RETURNING id
            """,
            workspace_id,
            folder_id,
            name,
            html,
            user_id,
        )
    except asyncpg.UniqueViolationError:
        row = await pool.fetchrow(
            """
            INSERT INTO pages (
                workspace_id, folder_id, name, content_html,
                content_type, html_layout, created_by
            ) VALUES ($1, $2, $3 || ' (imported)', $4, 'html', 'fixed-aspect', $5)
            RETURNING id
            """,
            workspace_id,
            folder_id,
            name,
            html,
            user_id,
        )

    return {
        "kind": "page",
        "page_id": str(row["id"]),
        "name": name,
        "slide_count": len(png_paths),
    }


async def _drive_download(client: httpx.AsyncClient, file_id: str) -> bytes:
    resp = await client.get(DRIVE_DOWNLOAD_URL.format(file_id=file_id))
    if resp.status_code == 404:
        raise RuntimeError("Drive file not found or not accessible to this account")
    resp.raise_for_status()
    return resp.content


async def _run(cmd: list[str], cwd: Path, timeout: int) -> None:
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        cwd=str(cwd),
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.STDOUT,
    )
    try:
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=timeout)
    except TimeoutError:
        proc.kill()
        await proc.wait()
        raise RuntimeError(f"{cmd[0]} timed out after {timeout}s")
    if proc.returncode != 0:
        tail = (stdout or b"").decode("utf-8", errors="replace")[-500:]
        raise RuntimeError(f"{cmd[0]} failed (exit {proc.returncode}): {tail}")


async def _convert_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    if shutil.which("libreoffice") is None and shutil.which("soffice") is None:
        raise RuntimeError(
            "libreoffice is not installed in the worker image — PPTX import requires it"
        )
    bin_name = "soffice" if shutil.which("soffice") else "libreoffice"
    await _run(
        [bin_name, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
        cwd=out_dir,
        timeout=LIBREOFFICE_TIMEOUT,
    )
    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        raise RuntimeError("libreoffice produced no PDF output")
    return pdf_path


async def _render_pdf_to_pngs(pdf_path: Path, out_dir: Path) -> list[Path]:
    if shutil.which("pdftoppm") is None:
        raise RuntimeError("pdftoppm (poppler-utils) is not installed in the worker image")
    prefix = out_dir / "slide"
    await _run(
        ["pdftoppm", "-r", str(PDFTOPPM_DPI), "-png", str(pdf_path), str(prefix)],
        cwd=out_dir,
        timeout=LIBREOFFICE_TIMEOUT,
    )
    pngs = sorted(out_dir.glob("slide-*.png"))
    if not pngs:
        raise RuntimeError("pdftoppm produced no PNGs")
    return pngs


def _extract_slide_text(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        runs: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    if t:
                        runs.append(t)
        texts.append("\n".join(runs).strip())
    return texts


def png_paths_to_data_uris(paths: list[Path]) -> list[str]:
    out: list[str] = []
    for p in paths:
        b64 = base64.b64encode(p.read_bytes()).decode("ascii")
        out.append(f"data:image/png;base64,{b64}")
    return out


def _build_slide_html(deck_name: str, image_uris: list[str], slide_texts: list[str]) -> str:
    """Build a fixed-aspect HTML page matching the slide-deck convention.

    Each `<section class="slide">` carries a background image and an
    invisible (visually hidden but semantically present) text overlay so
    the existing extraction/embedding pipeline can index the deck.
    """
    safe_name = html_mod.escape(deck_name)
    sections: list[str] = []
    for i, uri in enumerate(image_uris):
        text = slide_texts[i] if i < len(slide_texts) else ""
        safe_text = html_mod.escape(text)
        sections.append(
            '<section class="slide">'
            f'<img src="{uri}" alt="slide {i + 1}" />'
            '<div class="slide-text" aria-hidden="false">'
            f"{safe_text}"
            "</div>"
            "</section>"
        )
    body = "\n".join(sections)
    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8" />
<title>{safe_name}</title>
<style>
  html, body {{ margin: 0; padding: 0; background: #000; }}
  body > section.slide {{
    position: relative;
    width: 100%;
    aspect-ratio: 16 / 9;
    overflow: hidden;
  }}
  body > section.slide > img {{
    width: 100%;
    height: 100%;
    object-fit: contain;
    display: block;
  }}
  /* Visually-hidden but selectable + indexable. */
  body > section.slide > .slide-text {{
    position: absolute;
    width: 1px;
    height: 1px;
    overflow: hidden;
    clip: rect(0 0 0 0);
    white-space: pre-wrap;
  }}
</style>
</head>
<body>
{body}
</body>
</html>"""
