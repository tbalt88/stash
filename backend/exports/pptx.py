"""PPTX export — rasterized slides + invisible text overlay.

Uses Playwright to screenshot each `<section class="slide">` as a PNG
at 2x device scale, then assembles a 16:9 deck with one image per
slide and a transparent text layer harvested from the rendered DOM.
The text layer makes the exported PPTX selectable, copyable, and
searchable in PowerPoint / Keynote / Google Slides, without trying to
reconstruct native shapes (which would be a rabbit hole).

Reused by the Google Slides exporter (uploads the same PPTX to Drive
with `convertTo=true`).
"""

from __future__ import annotations

import io
import logging
import re
from uuid import UUID, uuid4

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from ..celery_app import celery
from ..database import get_pool
from ..services import storage_service
from ..tasks._celery_helpers import run_async
from .constants import (
    EXPORT_DEVICE_SCALE_FACTOR,
    SLIDE_HEIGHT_EMU,
    SLIDE_HEIGHT_PX,
    SLIDE_WIDTH_EMU,
    SLIDE_WIDTH_PX,
)

logger = logging.getLogger(__name__)

SECTION_RE = re.compile(
    r"<section\b[^>]*\bclass\s*=\s*[\"'][^\"']*\bslide\b[^\"']*[\"'][^>]*>",
    re.IGNORECASE,
)

# Text elements we harvest into the invisible PPTX text overlay. Limited
# to leaf-ish content elements — picking up <div> would double-count text
# because divs commonly wrap other elements we already capture.
TEXT_SELECTOR = "h1, h2, h3, h4, h5, h6, p, li, td, th, blockquote, figcaption"


def _count_slides(html: str) -> int:
    return max(1, len(SECTION_RE.findall(html or "")))


def _build_single_slide_html(source_html: str, slide_index: int) -> str:
    """Return HTML showing only the Nth <section class="slide">."""
    script = (
        "<script>(function(){var s=document.querySelectorAll('body > section.slide');"
        + f"var i={slide_index};"
        + "for(var k=0;k<s.length;k++){s[k].style.display=(k===i)?'':'none';}})();</script>"
    )
    if re.search(r"</body\s*>", source_html, flags=re.I):
        return re.sub(r"</body\s*>", script + "</body>", source_html, count=1, flags=re.I)
    return source_html + script


# JS executed inside the Playwright page to collect (text, bounds, size)
# tuples for the active slide. Bounds are in CSS px relative to the
# slide section so we can map them into PPTX EMUs by ratio.
_TEXT_HARVEST_JS = """
() => {
  const slide = document.querySelector('body > section.slide:not([style*="display: none"])')
            || document.querySelector('body > section.slide');
  if (!slide) return [];
  const slideRect = slide.getBoundingClientRect();
  const sel = "{SEL}";
  const out = [];
  slide.querySelectorAll(sel).forEach(el => {
    const text = el.innerText.trim();
    if (!text) return;
    const r = el.getBoundingClientRect();
    if (r.width < 4 || r.height < 4) return;
    const style = getComputedStyle(el);
    out.push({
      text,
      x: r.left - slideRect.left,
      y: r.top - slideRect.top,
      w: r.width,
      h: r.height,
      fontSize: parseFloat(style.fontSize) || 16,
      bold: parseInt(style.fontWeight, 10) >= 600,
      italic: style.fontStyle === 'italic',
    });
  });
  return out;
}
""".replace("{SEL}", TEXT_SELECTOR)


async def _capture_slide(page, slide_index: int, html: str) -> tuple[bytes, list[dict]]:
    """Render one slide, return (png_bytes, text_blocks).

    `text_blocks` are CSS-px coordinates relative to the slide section's
    top-left; the caller maps them to EMUs.
    """
    slide_html = _build_single_slide_html(html, slide_index)
    await page.set_content(slide_html, wait_until="networkidle")
    png = await page.screenshot(
        type="png",
        full_page=False,
        clip={"x": 0, "y": 0, "width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
    )
    try:
        text_blocks = await page.evaluate(_TEXT_HARVEST_JS)
    except Exception:
        logger.warning(
            "text harvest failed for slide %s; export will lack selectable text", slide_index
        )
        text_blocks = []
    return png, text_blocks


async def _capture_all_slides(html: str, count: int) -> list[tuple[bytes, list[dict]]]:
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        try:
            results: list[tuple[bytes, list[dict]]] = []
            context = await browser.new_context(
                viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
                device_scale_factor=EXPORT_DEVICE_SCALE_FACTOR,
            )
            try:
                for i in range(count):
                    page = await context.new_page()
                    try:
                        results.append(await _capture_slide(page, i, html))
                    finally:
                        await page.close()
            finally:
                await context.close()
            return results
        finally:
            await browser.close()


def _px_to_emu(px: float, total_px: int, total_emu: int) -> int:
    """Convert CSS px in a `total_px`-wide canvas to EMU in a
    `total_emu`-wide slide. Clamps to slide bounds so a stray
    out-of-bounds element doesn't produce a negative offset."""
    if total_px <= 0:
        return 0
    emu = int(px * total_emu / total_px)
    return max(0, min(total_emu, emu))


def _add_invisible_text(slide, text_blocks: list[dict]) -> None:
    """Add transparent text boxes positioned over the slide image so the
    text is selectable, copyable, and searchable in the export viewers."""
    for block in text_blocks:
        left = _px_to_emu(block["x"], SLIDE_WIDTH_PX, int(SLIDE_WIDTH_EMU))
        top = _px_to_emu(block["y"], SLIDE_HEIGHT_PX, int(SLIDE_HEIGHT_EMU))
        width = _px_to_emu(block["w"], SLIDE_WIDTH_PX, int(SLIDE_WIDTH_EMU))
        height = _px_to_emu(block["h"], SLIDE_HEIGHT_PX, int(SLIDE_HEIGHT_EMU))
        if width <= 0 or height <= 0:
            continue
        box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
        tf = box.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = block["text"]
        font = run.font
        # CSS px ≈ Pt * 96/72; the rendered image already encodes the
        # visual size, so this is purely for selection-rectangle sizing.
        font.size = Pt(max(6, int(block.get("fontSize", 16) * 72 / 96)))
        font.bold = bool(block.get("bold"))
        font.italic = bool(block.get("italic"))
        # Make the text invisible: matching white fill is fragile across
        # themes, so we drop alpha to 0 via the underlying XML. python-pptx
        # doesn't expose alpha directly on color, but the run's solid fill
        # accepts an `<a:alpha>` child via XML.
        font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_run_alpha(run, 0)


def _set_run_alpha(run, alpha_pct: int) -> None:
    """Set the alpha channel of a run's font color via the OOXML
    `<a:alpha val="..."/>` child. `alpha_pct` is in thousandths of a
    percent: 0 = fully transparent, 100000 = opaque."""
    from pptx.oxml.ns import qn

    rPr = run._r.get_or_add_rPr()
    solid = rPr.find(qn("a:solidFill"))
    if solid is None:
        # Color hasn't been written yet — touch it so python-pptx emits
        # the <a:solidFill> element, then re-find.
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        solid = rPr.find(qn("a:solidFill"))
    srgb = solid.find(qn("a:srgbClr")) if solid is not None else None
    if srgb is None:
        return
    # Drop any existing alpha child, then add a fresh one.
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    from lxml import etree

    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(alpha_pct * 1000)))


def _build_pptx(captures: list[tuple[bytes, list[dict]]]) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU
    blank = prs.slide_layouts[6]  # blank layout
    for png, text_blocks in captures:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(png),
            left=0,
            top=0,
            width=SLIDE_WIDTH_EMU,
            height=SLIDE_HEIGHT_EMU,
        )
        _add_invisible_text(slide, text_blocks)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def build_pptx_bytes_for_page(page_id: UUID) -> tuple[str, bytes]:
    """Render a page's slides to a PPTX. Returns (filename_stem, bytes).

    Exposed so the Google Slides exporter can call this without
    re-implementing the rendering pipeline.
    """
    pool = get_pool()
    page_row = await pool.fetchrow(
        "SELECT id, workspace_id, name, content_html, content_type, html_layout "
        "FROM pages WHERE id = $1",
        page_id,
    )
    if not page_row:
        raise RuntimeError("page not found")
    if page_row["content_type"] != "html" or page_row["html_layout"] != "fixed-aspect":
        raise RuntimeError("export requires a fixed-aspect HTML page")

    source_html = page_row["content_html"] or ""
    count = _count_slides(source_html)
    captures = await _capture_all_slides(source_html, count)
    pptx_bytes = _build_pptx(captures)
    return page_row["name"] or "slides", pptx_bytes


async def _export(user_id: UUID, page_id: UUID) -> dict:
    pool = get_pool()
    workspace_id = await pool.fetchval("SELECT workspace_id FROM pages WHERE id = $1", page_id)
    stem, pptx_bytes = await build_pptx_bytes_for_page(page_id)
    # SigV4 signing breaks on spaces / parens / other punctuation in S3 keys
    # once the URL gets percent-encoded. Collapse to an alnum-safe stem.
    safe_stem = re.sub(r"[^\w.-]+", "_", stem).strip("_") or "slides"
    filename = f"{safe_stem}-{uuid4().hex[:8]}.pptx"
    storage_key = await storage_service.upload_file(
        workspace_id=workspace_id,
        filename=filename,
        content=pptx_bytes,
        content_type=("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    )
    download_url = await storage_service.get_file_url(storage_key, expires_in=3600)
    return {
        "format": "pptx",
        "storage_key": storage_key,
        "download_url": download_url,
        "size_bytes": len(pptx_bytes),
    }


@celery.task(name="backend.exports.pptx.export_pptx")
def export_pptx(user_id: str, page_id: str) -> dict:
    return run_async(_export(UUID(user_id), UUID(page_id)))
